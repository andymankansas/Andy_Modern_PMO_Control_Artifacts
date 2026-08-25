"""Walking Deck design system and slide-block library.

Each block is a small renderer keyed by a stable id. `build_deck.py` reads a
config (brand + ordered list of selected blocks) and a content file, then calls
the matching renderer for each selected block. Blocks are individually optional
and degrade gracefully when content fields are missing, so any project can pick
only the slides it needs.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Callable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Fixed status palette (RAG) - not brandable so status reads consistently.
GREEN = RGBColor(16, 124, 16)
AMBER = RGBColor(242, 200, 17)
RED = RGBColor(164, 38, 44)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(48, 48, 48)
MID_GREY = RGBColor(105, 105, 105)
LIGHT_GREY = RGBColor(243, 245, 247)
LINE_GREY = RGBColor(214, 219, 224)

DEFAULT_BRAND = {
    "primary": "#0067B8",
    "accent": "#00A4EF",
    "dark": "#0F2337",
    "pale": "#D3F1FC",
    "light": "#ECF7FC",
    "logo_path": None,
    "hero_path": None,
    "title_font": "Arial",
    "body_font": "Segoe UI",
}


def _hex(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


@dataclass
class Brand:
    primary: RGBColor
    accent: RGBColor
    dark: RGBColor
    pale: RGBColor
    light: RGBColor
    title_font: str
    body_font: str
    logo_path: str | None = None
    hero_path: str | None = None

    @classmethod
    def from_dict(cls, data: dict | None) -> "Brand":
        merged = {**DEFAULT_BRAND, **(data or {})}
        return cls(
            primary=_hex(merged["primary"]),
            accent=_hex(merged["accent"]),
            dark=_hex(merged["dark"]),
            pale=_hex(merged["pale"]),
            light=_hex(merged["light"]),
            title_font=merged["title_font"],
            body_font=merged["body_font"],
            logo_path=merged["logo_path"],
            hero_path=merged["hero_path"],
        )

    def color(self, name: str | None) -> RGBColor:
        table = {
            "primary": self.primary, "accent": self.accent, "dark": self.dark,
            "navy": self.dark, "pale": self.pale, "light": self.light,
            "green": GREEN, "amber": AMBER, "red": RED,
            "white": WHITE, "charcoal": CHARCOAL, "grey": MID_GREY,
        }
        return table.get((name or "primary").lower(), self.primary)


@dataclass
class Ctx:
    brand: Brand
    page: int = 0
    footer: str = ""
    body_font: str = "Segoe UI"
    title_font: str = "Arial"


# --- primitive helpers -----------------------------------------------------

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, radius=False, line=None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, color)
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def add_chevron(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, color)
    return shape


def add_text(slide, x, y, w, h, text, size=12, color=CHARCOAL, bold=False,
             font="Segoe UI", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich(slide, x, y, w, h, lines, size=11, color=CHARCOAL, bullet_color=None,
             gap=3, margin=0.12, font="Segoe UI"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    for i, item in enumerate(lines):
        lead, body = (item + ["", ""])[:2] if isinstance(item, list) else (item, "")
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_after = Pt(gap)
        marker = p.add_run()
        marker.text = "\u203a "
        marker.font.name = font
        marker.font.size = Pt(size + 1)
        marker.font.bold = True
        marker.font.color.rgb = bullet_color or color
        if lead:
            r = p.add_run()
            r.text = str(lead)
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = color
        if body:
            r = p.add_run()
            r.text = str(body)
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


def add_title(slide, ctx: Ctx, title, takeaway=None):
    b = ctx.brand
    add_text(slide, 0.65, 0.42, 12.0, 0.5, title, 20, b.dark, True, font=ctx.title_font)
    add_rect(slide, 0.65, 0.98, 0.72, 0.05, b.accent)
    if takeaway:
        add_text(slide, 1.55, 0.86, 10.9, 0.30, takeaway, 10.5, b.primary, True, font=ctx.body_font)
    if ctx.footer:
        add_text(slide, 0.30, 7.08, 4.8, 0.24, ctx.footer, 6.5, MID_GREY, font=ctx.body_font)
    add_text(slide, 12.35, 7.08, 0.45, 0.22, str(ctx.page), 7, MID_GREY, align=PP_ALIGN.RIGHT, font=ctx.body_font)


def add_source(slide, ctx: Ctx, text):
    if text:
        add_text(slide, 6.25, 7.06, 5.75, 0.24, text, 5.8, MID_GREY, align=PP_ALIGN.RIGHT, margin=0, font=ctx.body_font)


def add_card(slide, ctx: Ctx, x, y, w, h, heading, body, accent, hsize=10.5, bsize=9.4):
    b = ctx.brand
    add_rect(slide, x, y, w, h, LIGHT_GREY, line=LINE_GREY)
    add_rect(slide, x, y, w, 0.10, accent)
    add_text(slide, x + 0.16, y + 0.18, w - 0.32, 0.42, heading, hsize, b.dark, True, font=ctx.title_font)
    if isinstance(body, str):
        add_text(slide, x + 0.16, y + 0.70, w - 0.32, h - 0.84, body, bsize, CHARCOAL, font=ctx.body_font)
    else:
        add_rich(slide, x + 0.08, y + 0.64, w - 0.16, h - 0.74, body, bsize, CHARCOAL, accent, gap=5, font=ctx.body_font)


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


# --- block renderers -------------------------------------------------------

def block_cover(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    if b.hero_path and Path(b.hero_path).exists():
        slide.shapes.add_picture(b.hero_path, 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        add_rect(slide, 0, 0, 13.333, 7.5, b.dark)
    add_rect(slide, 1.0, 1.05, 11.3, 5.25, b.primary)
    if c.get("eyebrow"):
        add_text(slide, 3.55, 1.95, 6.25, 0.38, c["eyebrow"], 9, WHITE, True, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_text(slide, 2.40, 2.55, 8.55, 1.48, c.get("title", "PROJECT\nWALKING DECK"), 28, WHITE, True,
             font=ctx.title_font, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if c.get("subtitle"):
        add_text(slide, 2.90, 4.34, 7.55, 0.50, c["subtitle"], 12, WHITE, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_rect(slide, 5.25, 5.05, 2.85, 0.03, WHITE)
    if c.get("stamp"):
        add_text(slide, 3.35, 5.20, 6.65, 0.32, c["stamp"], 9, WHITE, align=PP_ALIGN.CENTER, font=ctx.body_font)
    if b.logo_path and Path(b.logo_path).exists():
        slide.shapes.add_picture(b.logo_path, Inches(11.15), Inches(0.35), height=Inches(0.55))
    _notes(slide, c.get("notes"))


def block_agenda(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "AGENDA"), c.get("takeaway"))
    items = c.get("items", [])[:6]
    for i, it in enumerate(items):
        y = 1.45 + i * (5.2 / max(len(items), 1))
        accent = b.primary if i < max(1, len(items) - 2) else b.accent
        add_rect(slide, 0.80, y, 0.72, 0.72, accent)
        add_text(slide, 0.80, y, 0.72, 0.72, it.get("num", str(i + 1)), 16, WHITE, True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font=ctx.title_font)
        add_text(slide, 1.78, y + 0.02, 3.6, 0.30, it.get("head", ""), 11, b.dark, True, font=ctx.title_font)
        add_text(slide, 1.78, y + 0.36, 8.5, 0.30, it.get("body", ""), 10, CHARCOAL, font=ctx.body_font)
        add_rect(slide, 10.55, y + 0.34, 1.65, 0.04, LINE_GREY)
    _notes(slide, c.get("notes"))


def _cards_row(slide, ctx: Ctx, cards, top=2.34, height=4.02):
    b = ctx.brand
    n = min(len(cards), 4)
    if n == 0:
        return
    xs = {1: [4.9], 2: [2.9, 7.55], 3: [0.66, 4.55, 8.44], 4: [0.60, 3.66, 6.72, 9.78]}[n]
    w = {1: 3.55, 2: 4.9, 3: 3.55, 4: 2.85}[n]
    for x, card in zip(xs, cards[:n]):
        add_card(slide, ctx, x, top, w, height, card.get("head", ""), card.get("bullets") or card.get("body", ""),
                 b.color(card.get("accent")))


def block_problem(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "WHY WE'RE HERE"), c.get("takeaway"))
    ctx_band = c.get("context")
    top = 2.34
    if ctx_band:
        add_rect(slide, 0.66, 1.28, 11.48, 0.86, b.pale)
        add_text(slide, 0.86, 1.36, 2.4, 0.68, ctx_band.get("label", "WHO WE SERVE"), 10.5, b.primary, True,
                 font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, 3.05, 1.34, 9.0, 0.74, ctx_band.get("text", ""), 10, b.dark, valign=MSO_ANCHOR.MIDDLE, font=ctx.body_font)
    _cards_row(slide, ctx, c.get("cards", []), top=top, height=4.02)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


def block_vision(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "THE OPERATING MODEL"), c.get("takeaway"))
    if c.get("label"):
        add_rect(slide, 0.66, 1.30, 2.4, 0.36, b.dark)
        add_text(slide, 0.66, 1.30, 2.4, 0.36, c["label"], 10, WHITE, True, font=ctx.title_font,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    if c.get("statement"):
        add_text(slide, 3.20, 1.30, 9.0, 0.36, c["statement"], 9.5, CHARCOAL, valign=MSO_ANCHOR.MIDDLE, font=ctx.body_font)
    steps = c.get("steps", [])[:4]
    if steps:
        cw, gap = 2.55, 0.55
        x0 = (13.333 - (len(steps) * cw + (len(steps) - 1) * gap)) / 2
        for i, st in enumerate(steps):
            x = x0 + i * (cw + gap)
            accent = b.primary if i < len(steps) - 1 else b.accent
            add_rect(slide, x, 2.55, cw, 2.75, WHITE, line=LINE_GREY)
            add_rect(slide, x, 2.55, cw, 0.62, accent)
            add_text(slide, x, 2.55, cw, 0.62, st.get("num", str(i + 1)), 17, WHITE, True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0, font=ctx.title_font)
            add_text(slide, x + 0.14, 3.30, cw - 0.28, 0.66, st.get("head", ""), 11, b.dark, True,
                     font=ctx.title_font, align=PP_ALIGN.CENTER)
            add_text(slide, x + 0.14, 4.02, cw - 0.28, 1.18, st.get("body", ""), 9.4, CHARCOAL,
                     align=PP_ALIGN.CENTER, font=ctx.body_font)
            if i < len(steps) - 1:
                add_chevron(slide, x + cw + 0.10, 3.66, 0.34, 0.44, b.accent)
    if c.get("sowhat"):
        add_rect(slide, 0.74, 5.72, 11.85, 0.86, b.light)
        add_text(slide, 0.98, 5.80, 11.4, 0.70, c["sowhat"], 11.5, b.dark, True, valign=MSO_ANCHOR.MIDDLE, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_how_it_works(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "HOW IT WORKS"), c.get("takeaway"))
    stages = c.get("stages", [])[:9]
    accents = [b.primary, b.accent, GREEN]
    for i, st in enumerate(stages):
        col, row = i % 3, i // 3
        x = 0.72 + col * 4.05
        y = 1.42 + row * 1.42
        accent = accents[row % 3]
        add_rect(slide, x, y, 3.55, 1.14, LIGHT_GREY, line=LINE_GREY)
        add_rect(slide, x, y, 0.52, 1.14, accent)
        add_text(slide, x, y, 0.52, 1.14, st.get("num", str(i + 1)), 16, WHITE, True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0, font=ctx.title_font)
        add_text(slide, x + 0.66, y + 0.14, 2.75, 0.34, st.get("head", ""), 9.6, b.dark, True, font=ctx.title_font)
        add_text(slide, x + 0.66, y + 0.56, 2.75, 0.46, st.get("body", ""), 8.6, CHARCOAL, font=ctx.body_font)
        if col < 2 and i + 1 < len(stages):
            add_chevron(slide, x + 3.62, y + 0.40, 0.28, 0.34, LINE_GREY)
    if c.get("governance"):
        add_rect(slide, 0.72, 5.86, 11.87, 0.80, b.pale)
        add_text(slide, 0.96, 5.90, 2.1, 0.72, c.get("governance_label", "WHY IT\nMATTERS"), 9.5, b.primary, True,
                 font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, 3.05, 5.90, 9.35, 0.72, c["governance"], 9.6, b.dark, valign=MSO_ANCHOR.MIDDLE, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_roles(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "EACH PLATFORM HAS ONE CLEAR JOB"), c.get("takeaway"))
    cards = c.get("cards", [])[:4]
    n = max(len(cards), 1)
    xs = {1: [5.0], 2: [2.4, 7.4], 3: [0.66, 4.55, 8.44], 4: [0.60, 3.66, 6.72, 9.78]}[n]
    cw = {1: 3.3, 2: 3.5, 3: 3.55, 4: 2.85}[n]
    for x, card in zip(xs, cards):
        accent = b.color(card.get("accent"))
        add_rect(slide, x, 1.50, cw, 4.95, WHITE, line=LINE_GREY)
        add_rect(slide, x, 1.50, cw, 0.82, accent)
        add_text(slide, x + 0.10, 1.60, cw - 0.20, 0.30, card.get("name", ""), 12, WHITE, True,
                 font=ctx.title_font, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, x + 0.10, 1.98, cw - 0.20, 0.24, card.get("role", ""), 8, WHITE, True,
                 align=PP_ALIGN.CENTER, margin=0, font=ctx.body_font)
        add_rich(slide, x + 0.10, 2.52, cw - 0.18, 3.80, card.get("rows", []), 9.2, CHARCOAL, accent, gap=8, font=ctx.body_font)
    if c.get("footnote"):
        add_text(slide, 0.60, 6.60, 12.1, 0.30, c["footnote"], 8.6, MID_GREY, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


def block_scope(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "SCOPE"), c.get("takeaway"))
    add_rect(slide, 0.66, 1.45, 5.75, 5.0, WHITE, line=LINE_GREY)
    add_rect(slide, 0.66, 1.45, 5.75, 0.55, b.primary)
    add_text(slide, 0.86, 1.45, 5.4, 0.55, "IN SCOPE", 12, WHITE, True, font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE)
    add_rich(slide, 0.80, 2.15, 5.5, 4.1, [[i, ""] if isinstance(i, str) else i for i in c.get("in_scope", [])],
             10, CHARCOAL, b.primary, gap=7, font=ctx.body_font)
    add_rect(slide, 6.92, 1.45, 5.75, 5.0, WHITE, line=LINE_GREY)
    add_rect(slide, 6.92, 1.45, 5.75, 0.55, MID_GREY)
    add_text(slide, 7.12, 1.45, 5.4, 0.55, "OUT OF SCOPE", 12, WHITE, True, font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE)
    add_rich(slide, 7.06, 2.15, 5.5, 4.1, [[i, ""] if isinstance(i, str) else i for i in c.get("out_scope", [])],
             10, CHARCOAL, MID_GREY, gap=7, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_roadmap(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "THE ROADMAP"), c.get("takeaway"))
    waves = c.get("waves", [])[:3]
    xs = [0.68, 4.98, 9.28]
    for x, w in zip(xs, waves):
        accent = b.color(w.get("accent"))
        add_rect(slide, x, 1.48, 3.38, 1.66, WHITE, line=LINE_GREY)
        add_rect(slide, x, 1.48, 3.38, 0.13, accent)
        add_text(slide, x + 0.16, 1.72, 2.15, 0.30, w.get("head", ""), 10.5, b.dark, True, font=ctx.title_font)
        add_text(slide, x + 2.30, 1.72, 0.92, 0.30, w.get("status", ""), 7.6, accent, True, align=PP_ALIGN.RIGHT, margin=0, font=ctx.body_font)
        add_text(slide, x + 0.16, 2.12, 3.05, 0.26, w.get("sub", ""), 8.2, accent, True, font=ctx.body_font)
        add_text(slide, x + 0.16, 2.46, 3.05, 0.60, w.get("body", ""), 8.3, CHARCOAL, font=ctx.body_font)
    if len(waves) == 3:
        for x in (4.20, 8.50):
            add_chevron(slide, x, 2.05, 0.42, 0.48, b.accent)
    cycle = c.get("cycle", [])[:6]
    if cycle:
        add_text(slide, 0.72, 3.36, 9.0, 0.30, c.get("cycle_label", "EVERY WAVE RUNS THE SAME CYCLE"), 10.5, b.dark, True, font=ctx.title_font)
        for i, st in enumerate(cycle):
            x = 0.72 + i * 2.03
            accent = b.primary if i < len(cycle) - 1 else GREEN
            add_rect(slide, x, 3.90, 1.73, 1.30, LIGHT_GREY, line=LINE_GREY)
            add_rect(slide, x, 3.90, 0.42, 1.30, accent)
            add_text(slide, x, 3.90, 0.42, 1.30, st.get("num", str(i + 1)), 12, WHITE, True,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0, font=ctx.title_font)
            add_text(slide, x + 0.52, 4.08, 1.12, 0.36, st.get("head", ""), 8.0, b.dark, True, font=ctx.title_font)
            add_text(slide, x + 0.52, 4.56, 1.14, 0.42, st.get("body", ""), 7.2, CHARCOAL, font=ctx.body_font)
            if i < len(cycle) - 1:
                add_chevron(slide, x + 1.76, 4.38, 0.22, 0.32, LINE_GREY)
    if c.get("outcome"):
        add_rect(slide, 0.72, 5.66, 11.87, 1.00, b.pale)
        add_text(slide, 0.96, 5.84, 2.2, 0.30, c.get("outcome_label", "OUTCOME"), 9.5, b.primary, True, font=ctx.title_font)
        add_text(slide, 3.15, 5.76, 9.2, 0.80, c["outcome"], 10.5, b.dark, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_milestones(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "DELIVERY MILESTONES"), c.get("takeaway"))
    items = c.get("items", [])[:7]
    status_color = {"complete": GREEN, "active": AMBER, "upcoming": b.primary}
    add_rect(slide, 0.86, 3.30, 11.36, 0.08, LINE_GREY)
    for i, m in enumerate(items):
        x = 0.66 + i * (11.5 / max(len(items), 1))
        accent = status_color.get(m.get("status", "upcoming"), b.primary)
        add_text(slide, x, 1.55, 1.54, 0.30, m.get("date", ""), 9.5, accent, True, font=ctx.title_font, align=PP_ALIGN.CENTER)
        add_text(slide, x, 2.00, 1.54, 0.90, m.get("head", ""), 9.5, b.dark, True, font=ctx.title_font, align=PP_ALIGN.CENTER)
        add_rect(slide, x + 0.68, 3.06, 0.18, 0.56, accent)
        add_text(slide, x, 3.78, 1.54, 1.16, m.get("body", ""), 8.2, CHARCOAL, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_rect(slide, 0.78, 5.35, 11.44, 0.52, LIGHT_GREY, line=LINE_GREY)
    for i, (label, color) in enumerate([("Complete", GREEN), ("In progress", AMBER), ("Upcoming", b.primary)]):
        bx = 3.9 + i * 1.9
        add_rect(slide, bx, 5.50, 0.22, 0.22, color)
        add_text(slide, bx + 0.28, 5.46, 1.5, 0.30, label, 9, CHARCOAL, font=ctx.body_font)
    if c.get("critical_path"):
        add_rect(slide, 0.78, 6.06, 11.44, 0.80, b.pale)
        add_text(slide, 1.02, 6.24, 11.0, 0.44, c["critical_path"], 10.5, b.dark, True, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


def block_current_state(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "WHERE WE ARE NOW"), c.get("takeaway"))
    top = 2.30
    if c.get("mode"):
        add_rect(slide, 0.66, 1.30, 11.48, 0.66, AMBER)
        add_text(slide, 0.90, 1.44, 11.0, 0.34, c["mode"], 10.6, b.dark, True, align=PP_ALIGN.CENTER, font=ctx.body_font)
    _cards_row(slide, ctx, c.get("cards", []), top=top, height=4.05)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


def block_team(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "DELIVERY DEPENDS ON SHARED OWNERSHIP"), c.get("takeaway"))
    add_rect(slide, 0.66, 1.38, 7.28, 5.42, WHITE, line=LINE_GREY)
    add_rect(slide, 8.22, 1.38, 3.92, 5.42, b.light)
    add_text(slide, 0.92, 1.64, 6.78, 0.30, c.get("owners_label", "CORE OWNERSHIP MAP"), 11, b.dark, True, font=ctx.title_font)
    owners = c.get("owners", [])[:7]
    for i, o in enumerate(owners):
        y = 2.10 + i * 0.62
        add_rect(slide, 0.92, y, 0.10, 0.43, b.color(o.get("accent")))
        add_text(slide, 1.18, y - 0.04, 2.62, 0.30, o.get("role", ""), 7.8, b.dark, True, font=ctx.title_font)
        add_text(slide, 3.88, y - 0.01, 3.66, 0.38, o.get("names", ""), 8.8, CHARCOAL, font=ctx.body_font)
    add_text(slide, 8.55, 1.68, 3.28, 0.30, c.get("rhythm_label", "OPERATING RHYTHM"), 11, b.dark, True, font=ctx.title_font)
    add_rich(slide, 8.46, 2.13, 3.40, 3.98, c.get("rhythm", []), 9.2, CHARCOAL, b.primary, gap=7, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_next_steps(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "NEXT STEPS"), c.get("takeaway"))
    add_rect(slide, 0.66, 1.37, 7.42, 5.48, WHITE, line=LINE_GREY)
    add_rect(slide, 8.38, 1.37, 3.76, 5.48, b.pale)
    add_text(slide, 0.94, 1.68, 6.82, 0.34, c.get("moves_label", "NEXT MOVES"), 11, b.dark, True, font=ctx.title_font)
    moves = c.get("moves", [])[:6]
    for i, m in enumerate(moves):
        y = 2.16 + i * (4.4 / max(len(moves), 1))
        add_rect(slide, 0.95, y, 0.50, 0.50, b.primary if i < 2 else b.accent)
        add_text(slide, 0.95, y, 0.50, 0.50, m.get("num", str(i + 1)), 11, WHITE, True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0, font=ctx.title_font)
        add_text(slide, 1.65, y - 0.01, 4.70, 0.52, m.get("action", ""), 9.8, CHARCOAL, True, font=ctx.body_font)
        add_text(slide, 6.44, y + 0.04, 1.18, 0.26, m.get("timing", ""), 8.5, b.primary, True, align=PP_ALIGN.RIGHT, font=ctx.body_font)
    add_text(slide, 8.72, 1.72, 3.08, 0.34, c.get("help_label", "LEADERSHIP HELP WANTED"), 10.5, b.dark, True, font=ctx.title_font)
    add_rich(slide, 8.58, 2.20, 3.28, 3.74, c.get("help", []), 9.4, CHARCOAL, b.primary, gap=8, font=ctx.body_font)
    if c.get("success"):
        add_text(slide, 8.70, 6.16, 3.04, 0.44, c["success"], 8.8, b.primary, True, align=PP_ALIGN.CENTER, font=ctx.body_font)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


def block_closer(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_title(slide, ctx, c.get("title", "WHAT'S NEXT"), c.get("takeaway"))
    tiers = c.get("tiers", [])[:3]
    geo = [(0.85, 4.45, 1.70), (4.90, 3.65, 2.50), (8.95, 2.85, 3.30)]
    for (x, y, h), t in zip(geo, tiers):
        accent = b.color(t.get("accent"))
        add_rect(slide, x, y, 3.60, h, LIGHT_GREY, line=LINE_GREY)
        add_rect(slide, x, y, 3.60, 0.66, accent)
        add_text(slide, x + 0.16, y + 0.10, 3.30, 0.46, t.get("head", ""), 12, WHITE, True, font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE, margin=0)
        add_text(slide, x + 0.18, y + 0.80, 3.26, 0.28, t.get("when", ""), 9, accent, True, font=ctx.body_font)
        add_text(slide, x + 0.18, y + 1.16, 3.26, h - 1.30, t.get("body", ""), 9.6, CHARCOAL, font=ctx.body_font)
    if len(tiers) == 3:
        add_chevron(slide, 4.10, 2.95, 0.60, 0.60, LINE_GREY)
        add_chevron(slide, 8.15, 2.95, 0.60, 0.60, LINE_GREY)
    if c.get("banner"):
        add_rect(slide, 0.85, 6.36, 11.70, 0.56, b.dark)
        add_text(slide, 0.85, 6.36, 11.70, 0.56, c["banner"], 12.5, WHITE, True, font=ctx.title_font,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    _notes(slide, c.get("notes"))


def block_section(slide, ctx: Ctx, c: dict):
    b = ctx.brand
    add_rect(slide, 0, 0, 13.333, 7.5, b.dark)
    add_rect(slide, 0, 3.2, 13.333, 0.04, b.accent)
    add_text(slide, 1.0, 2.7, 11.3, 1.4, c.get("title", "SECTION"), 30, WHITE, True,
             font=ctx.title_font, valign=MSO_ANCHOR.MIDDLE)
    if c.get("subtitle"):
        add_text(slide, 1.02, 4.2, 11.3, 0.5, c["subtitle"], 13, b.pale, font=ctx.body_font)
    _notes(slide, c.get("notes"))


def block_bullets(slide, ctx: Ctx, c: dict):
    add_title(slide, ctx, c.get("title", "OVERVIEW"), c.get("takeaway"))
    add_rich(slide, 0.80, 1.55, 11.6, 5.0,
             [[i, ""] if isinstance(i, str) else i for i in c.get("bullets", [])],
             12, CHARCOAL, ctx.brand.primary, gap=8, font=ctx.body_font)
    add_source(slide, ctx, c.get("source"))
    _notes(slide, c.get("notes"))


REGISTRY: dict[str, Callable] = {
    "cover": block_cover,
    "agenda": block_agenda,
    "problem": block_problem,
    "vision": block_vision,
    "how_it_works": block_how_it_works,
    "roles": block_roles,
    "scope": block_scope,
    "roadmap": block_roadmap,
    "milestones": block_milestones,
    "current_state": block_current_state,
    "team": block_team,
    "next_steps": block_next_steps,
    "closer": block_closer,
    "section": block_section,
    "bullets": block_bullets,
}

BLOCK_IDS = list(REGISTRY.keys())
