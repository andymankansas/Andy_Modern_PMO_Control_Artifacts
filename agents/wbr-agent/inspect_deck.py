"""inspect_deck.py - dump a deck's structure so you can map it to wbr_config.json.

Cross-platform (uses python-pptx only, no PowerPoint install needed). Run this on
your own status template, then use the output to fill the shape anchors in
wbr_config.json. An LLM can do the mapping for you: give it this JSON plus
wbr_config.example.json and ask it to produce your wbr_config.json.

Usage (from the extracted package folder):
    python inspect_deck.py "path/to/your_template.pptx"
    python inspect_deck.py "path/to/your_template.pptx" --out deck-structure.json
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400


def _inches(value) -> float:
    return round((value or 0) / EMU_PER_INCH, 3)


def inspect(path: Path) -> dict:
    prs = Presentation(str(path))
    report = {
        "source": str(path),
        "slide_count": len(prs.slides),
        "width_inches": _inches(prs.slide_width),
        "height_inches": _inches(prs.slide_height),
        "slides": [],
    }
    for number, slide in enumerate(prs.slides, start=1):
        slide_report = {
            "number": number,
            "layout": slide.slide_layout.name,
            "title": slide.shapes.title.text.strip() if slide.shapes.title else "",
            "shapes": [],
        }
        for shape in slide.shapes:
            entry = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": _inches(shape.left),
                "top": _inches(shape.top),
                "width": _inches(shape.width),
                "height": _inches(shape.height),
                "text": shape.text.strip() if shape.has_text_frame else "",
                "is_group": shape.shape_type == MSO_SHAPE_TYPE.GROUP,
            }
            if shape.has_table:
                table = shape.table
                entry["table"] = {
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "cells": [[cell.text.strip() for cell in row.cells] for row in table.rows],
                }
            if shape.has_chart:
                entry["chart"] = {"type": str(shape.chart.chart_type)}
            slide_report["shapes"].append(entry)
        report["slides"].append(slide_report)
    return report


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(prog="inspect_deck")
    ap.add_argument("deck", help="Path to the .pptx template to inspect")
    ap.add_argument("--out", default=None, help="Write JSON here (default: print to stdout)")
    args = ap.parse_args(argv)

    path = Path(args.deck)
    if not path.exists():
        print(f"ERROR: deck not found: {path}")
        return 2
    report = inspect(path)
    text = json.dumps(report, indent=2, ensure_ascii=False)
    if args.out:
        Path(args.out).write_text(text, encoding="utf-8")
        print(f"Wrote {args.out} ({report['slide_count']} slides).")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
