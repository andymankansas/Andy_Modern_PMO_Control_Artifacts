"""wbr_lib.py - shared helpers for the WBR (Weekly Business Review) agent.

The WBR agent rolls the prior week's deck forward: it copies the previous
PPTX to a new dated file and surgically updates dates, RAG status, program
narratives, the risk register, the reporting slide, and support metrics.
Editorial content comes from a weekly input file; the Project Context Layer
knowledge base supplies a change digest that surfaces what moved week over week.
"""
from __future__ import annotations

import datetime as _dt
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Deck typography: the template uses an en-dash between range endpoints. Kept as
# an escape so this source file carries no literal em/en-dash characters.
EN_DASH = "\u2013"

HERE = Path(__file__).resolve().parent

# Lines whose whole text is one of these render as bold section headers.
_BODY_HEADERS = {
    "accomplishments", "in progress", "risks / blockers", "risks/blockers",
    "risks", "blockers", "next steps", "highlights", "focus",
}
_BULLET_STARTS = ("\u2022", "-", "*")


# --------------------------------------------------------------------------
# Config + weekly input
# --------------------------------------------------------------------------
def load_config(path: str | Path | None = None) -> dict:
    path = Path(path) if path else HERE / "wbr_config.json"
    return json.loads(path.read_text(encoding="utf-8"))


def load_weekly_input(path: str | Path) -> dict:
    return json.loads(Path(path).read_text(encoding="utf-8"))


# --------------------------------------------------------------------------
# Dates
# --------------------------------------------------------------------------
def parse_date(value: str) -> _dt.date:
    """Accept ISO (2026-08-31) or M/D/YYYY (8/31/2026)."""
    value = value.strip()
    for fmt in ("%Y-%m-%d", "%m/%d/%Y", "%m/%d/%y"):
        try:
            return _dt.datetime.strptime(value, fmt).date()
        except ValueError:
            continue
    raise ValueError(f"Unrecognized date: {value!r}")


def _no_pad(n: int) -> str:
    return str(int(n))


def week_forms(start: _dt.date, end: _dt.date) -> dict:
    """Every string form of a reporting week that appears in the deck."""
    same_month = start.month == end.month
    if same_month:
        long_form = f"{start.strftime('%B')} {start.day}-{end.day}, {end.year}"
    else:
        long_form = (f"{start.strftime('%B')} {start.day}-"
                     f"{end.strftime('%B')} {end.day}, {end.year}")
    slash = (f"{_no_pad(start.month)}/{_no_pad(start.day)}/{start.year} {EN_DASH} "
             f"{_no_pad(end.month)}/{_no_pad(end.day)}/{end.year}")
    week_of_long = f"Week of {long_form}"
    return {"long": long_form, "slash": slash, "week_of_long": week_of_long}


def metrics_forms(start: _dt.date, end: _dt.date) -> dict:
    """String forms of the support-metrics week (slide 9 + timeline updated)."""
    slash = (f"{start.month:02d}/{start.day:02d}/{start.year}- "
             f"{end.month:02d}/{end.day:02d}/{end.year}")
    week_start = f"Week of {_ordinal(start.day)} {start.strftime('%B')}"
    week_end = f"Week of {_ordinal(end.day)} {end.strftime('%B')}"
    updated = f"Updated {_no_pad(end.month)}/{_no_pad(end.day)}/{end.year}"
    return {"slash": slash, "week_start": week_start,
            "week_end": week_end, "updated": updated}


def _ordinal(day: int) -> str:
    if 11 <= day % 100 <= 13:
        suffix = "th"
    else:
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
    return f"{day}{suffix}"


def output_filename(pattern: str, start: _dt.date, end: _dt.date) -> str:
    """Fill {start_month} {start_day} {end_month} {end_day} {year} tokens."""
    same_month = start.month == end.month
    end_label = f"{end.day}" if same_month else f"{end.strftime('%B')} {end.day}"
    return pattern.format(
        start_month=start.strftime("%B"), start_day=start.day,
        end_month=end.strftime("%B"), end_day=end.day,
        end_label=end_label, year=end.year,
    )


# --------------------------------------------------------------------------
# Shape access
# --------------------------------------------------------------------------
def slide_at(prs: Presentation, number: int):
    return prs.slides[number - 1]


def find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _style_donor(text_frame):
    """Capture font attributes from the first run to reapply after a rewrite."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            color = None
            try:
                if font.color and font.color.type is not None:
                    color = RGBColor(*font.color.rgb)
            except Exception:
                color = None
            return {"name": font.name, "size": font.size,
                    "bold": font.bold, "italic": font.italic, "color": color}
    return {"name": None, "size": None, "bold": None, "italic": None, "color": None}


def set_text(shape, text: str) -> None:
    """Replace a shape's text, preserving the first run's look.

    Multi-line text (split on newline) is rebuilt as separate paragraphs so the
    deck's line breaks survive.
    """
    tf = shape.text_frame
    donor = _style_donor(tf)
    lines = str(text).split("\n")
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        _apply_font(run.font, donor)


def _apply_font(font, donor: dict) -> None:
    if donor.get("name"):
        font.name = donor["name"]
    if donor.get("size"):
        font.size = donor["size"]
    if donor.get("bold") is not None:
        font.bold = donor["bold"]
    if donor.get("italic") is not None:
        font.italic = donor["italic"]
    if donor.get("color") is not None:
        font.color.rgb = donor["color"]


def set_rich_text(shape, text: str) -> None:
    """Replace text while restoring the deck's bold-lead-in / header / bullet look.

    - A line matching a known section header renders bold.
    - A 'Label: rest' line renders the label bold and the remainder regular.
    - Bullet lines and plain lines render regular.
    Font name, size, and color are inherited from the shape's first run.
    """
    tf = shape.text_frame
    donor = _style_donor(tf)
    lines = str(text).split("\n")
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.strip()
        if not stripped:
            continue
        low = stripped.lower().rstrip(":")
        is_bullet = line.lstrip().startswith(_BULLET_STARTS)
        if not is_bullet and (low in _BODY_HEADERS or stripped.endswith(":")):
            run = para.add_run()
            run.text = line
            _apply_font(run.font, donor)
            run.font.bold = True
        elif not is_bullet and ": " in line:
            label, rest = line.split(": ", 1)
            r1 = para.add_run()
            r1.text = label + ": "
            _apply_font(r1.font, donor)
            r1.font.bold = True
            r2 = para.add_run()
            r2.text = rest
            _apply_font(r2.font, donor)
            r2.font.bold = False
        else:
            run = para.add_run()
            run.text = line
            _apply_font(run.font, donor)
            run.font.bold = False


def set_fill(shape, rgb_hex: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(rgb_hex)


def set_cell_text(cell, text: str) -> None:
    donor = _style_donor(cell.text_frame)
    cell.text_frame.clear()
    para = cell.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = str(text)
    _apply_font(run.font, donor)


def replace_text_everywhere(prs: Presentation, old: str, new: str) -> int:
    """Replace a literal string across every text run in the deck."""
    count = 0
    for slide in prs.slides:
        count += _replace_in_shapes(slide.shapes, old, new)
    return count


def _replace_in_shapes(shapes, old: str, new: str) -> int:
    count = 0
    for shape in shapes:
        if shape.shape_type == 6:  # group
            count += _replace_in_shapes(shape.shapes, old, new)
            continue
        if shape.has_text_frame:
            count += _replace_in_text_frame(shape.text_frame, old, new)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    count += _replace_in_text_frame(cell.text_frame, old, new)
    return count


def _replace_in_text_frame(tf, old: str, new: str) -> int:
    count = 0
    for para in tf.paragraphs:
        # Join runs to catch strings split across run boundaries, then rewrite
        # the first run and blank the rest only when a match is present.
        full = "".join(run.text for run in para.runs)
        if old in full:
            updated = full.replace(old, new)
            count += full.count(old)
            if para.runs:
                para.runs[0].text = updated
                for run in para.runs[1:]:
                    run.text = ""
    return count


def strip_highlights(prs: Presentation) -> None:
    for slide in prs.slides:
        _strip_in_shapes(slide.shapes)


def _strip_in_shapes(shapes) -> None:
    for shape in shapes:
        if shape.shape_type == 6:
            _strip_in_shapes(shape.shapes)
            continue
        frames = []
        if shape.has_text_frame:
            frames.append(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    frames.append(cell.text_frame)
        for tf in frames:
            for para in tf.paragraphs:
                for run in para.runs:
                    rpr = run._r.find(qn("a:rPr"))
                    if rpr is not None:
                        for hl in rpr.findall(qn("a:highlight")):
                            rpr.remove(hl)


# --------------------------------------------------------------------------
# Deck week detection
# --------------------------------------------------------------------------
_SLASH_RANGE = re.compile(
    r"(\d{1,2})/(\d{1,2})/(\d{4})\s*[\u2013\u2014-]\s*(\d{1,2})/(\d{1,2})/(\d{4})")


def detect_report_week(prs: Presentation, subtitle_slide: int, subtitle_shape: str):
    """Return (start, end) of the deck's current reporting week, or None."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            m = _SLASH_RANGE.search(shape.text_frame.text)
            if m:
                start = _dt.date(int(m.group(3)), int(m.group(1)), int(m.group(2)))
                end = _dt.date(int(m.group(6)), int(m.group(4)), int(m.group(5)))
                return start, end
    return None
